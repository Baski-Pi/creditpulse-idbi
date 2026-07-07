"""Fill the official IDBI Innovate template with CreditPulse content.

Output: docs/CreditPulse_Track4.pptx  (export to PDF once links are final)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path(__file__).resolve().parents[1]
DOCS = BASE / "docs-"
TEMPLATE = DOCS / "Prototype Submission Deck _ IDBI Innovate.pptx"
OUT = DOCS / "CreditPulse_Track4.pptx"

INK = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0xB9, 0x1C, 0x1C)

prs = Presentation(str(TEMPLATE))
SW = prs.slide_width
SH = prs.slide_height


def add_bullets(slide, items, top=1.6, left=0.6, width=None, height=None,
                size=15, line_gap=6):
    width = width or (SW / 914400 - 1.2)
    height = height or (SH / 914400 - top - 0.4)
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, lvl, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = lvl
        p.space_after = Pt(line_gap)
        for r in p.runs:
            r.font.size = Pt(size if lvl == 0 else size - 2)
            r.font.bold = bold
            r.font.color.rgb = INK
    return tb


def add_picture_fit(slide, img, top=1.5, bottom_margin=0.3):
    from PIL import Image
    with Image.open(img) as im:
        w, h = im.size
    max_w = SW / 914400 - 1.0
    max_h = SH / 914400 - top - bottom_margin
    scale = min(max_w / (w / 160), max_h / (h / 160))
    pw, ph = (w / 160) * scale, (h / 160) * scale
    left = (SW / 914400 - pw) / 2
    slide.shapes.add_picture(str(img), Inches(left), Inches(top),
                             Inches(pw), Inches(ph))


S = prs.slides

# ---- Slide 1: Team details
add_bullets(S[0], [
    ("Team name:  Credit Pulse", 0, True),
    ("Team leader:  Baskar B  (baskarbssb@gmail.com)", 0, True),
    ("Problem statement:  Track 4 - Default Prediction Model", 0, True),
    ("", 0, False),
    ("CreditPulse - knows WHEN, explains WHY, lets the officer decide.", 0, False),
], top=2.2, size=18)

# ---- Slide 2: Brief about the idea
add_bullets(S[1], [
    ("CreditPulse is a default-risk early-warning system for loan accounts already running in the bank.", 0, True),
    ("It predicts not just IF an account will default, but WHEN - a month-by-month default-probability curve up to 12 months ahead.", 0, False),
    ("Every account is re-scored monthly using the three pillars IDBI specified: borrower behavior, the bank's internal database, and public-domain data.", 0, False),
    ("Output is exactly what the bank asked for: Red / Amber / Green risk buckets - plus an intervention window ('risk peaks around month 8') and plain-language reason codes for every alert.", 0, False),
    ("Human-in-the-loop by design: CreditPulse flags and explains; the credit officer decides and acts.", 0, True),
], size=16)

# ---- Slide 3: Opportunities / differentiation / USP
add_bullets(S[2], [
    ("How is it different?", 0, True),
    ("Competing approaches answer 'will this account default - yes/no'. CreditPulse answers WHEN, WHY, and WHAT TO DO NOW.", 1, False),
    ("Time-to-default is modeled with a discrete-time hazard model - the same construction banks use for IFRS 9 lifetime PD - built here as a working prototype.", 1, False),
    ("Rule triggers deliberately mirror RBI's Early Warning Signals framework (Master Directions, July 2024) - regulator-aligned from day one.", 1, False),
    ("How does it solve the problem?", 0, True),
    ("A 12-month PD curve per account turns prediction into an actionable window: who to call, and how much time is left to act.", 1, False),
    ("Proven blind on a full hold-out year: accounts we bucketed RED actually defaulted at 11.0% vs 3.3% for GREEN (3.3x separation).", 1, False),
    ("USP", 0, True),
    ("The only watchlist that tells the credit officer how much time they have - and why.", 1, True),
], size=15)

# ---- Slide 4: Features
add_bullets(S[3], [
    ("RAG portfolio watchlist (the bank's requested output), sorted by 12-month default probability", 0, False),
    ("Account 360: PD curve, risk-peak month, account profile, and recommended actions with officer sign-off", 0, False),
    ("Plain-language reason codes for every alert (SHAP-based) - 'utilization 95% and rising', 'EMI burden high vs income'", 0, False),
    ("Transparent EWS rule triggers (RBI-style) working alongside the ML score - officers can verify them by hand", 0, False),
    ("Blind-validation view: the dashboard proves its own buckets on historical outcomes", 0, False),
    ("Calibrated probabilities (predicted 14.6% vs observed 14.2% in the riskiest decile) - PDs officers can trust", 0, False),
    ("Segment-aware scoring (vintage, purpose, thin-history) feeding one unified score - as IDBI described", 0, False),
    ("Monthly re-scoring of the entire book; stability monitored per quarter", 0, False),
], size=15)

# ---- Slide 5: Process flow
add_picture_fit(S[4], DOCS / "diagram_flow.png", top=1.5)

# ---- Slide 6: Wireframes (optional) - use watchlist screenshot
add_picture_fit(S[5], DOCS / "1.png", top=1.5)

# ---- Slide 7: Architecture
add_picture_fit(S[6], DOCS / "diagram_arch.png", top=1.4)

# ---- Slide 8: Technologies
add_bullets(S[7], [
    ("Modeling: Python 3.11, LightGBM (gradient boosting), scikit-learn, SHAP explainability, pandas/pyarrow", 0, False),
    ("Method: discrete-time hazard model (loan-month panel, 6.2M rows) for the 12-month PD term structure", 0, False),
    ("Dashboard: Streamlit; prototype deployed on Streamlit Cloud straight from the public GitHub repo", 0, False),
    ("Production (sandbox-ready): AWS S3 + Glue data lake, SageMaker scoring endpoint, Lambda rules engine, API Gateway", 0, False),
    ("Fully open-source stack - zero license cost, deployable on cloud or on-prem per bank policy", 0, False),
    ("Training data: 2.26M real loans with true default timing (Lending Club) + 30k-customer behavioral dataset (UCI); feature schema maps 1:1 to CBS transactions, CIBIL bureau fields and RBI public data", 0, False),
], size=15)

# ---- Slide 9: Cost
add_bullets(S[8], [
    ("Prototype: zero cost (open-source software, free-tier hosting)", 0, True),
    ("Estimated pilot run-cost on AWS (full-book monthly scoring):", 0, False),
    ("S3 + Glue monthly batch: ~USD 150 / month", 1, False),
    ("SageMaker scoring endpoint: ~USD 250 / month", 1, False),
    ("Lambda rules + API Gateway: < USD 50 / month", 1, False),
    ("Total under ~USD 500 / month; no software license fees at any scale", 0, True),
    ("On-prem option: same stack runs on commodity hardware (no GPU required - LightGBM is CPU-native)", 0, False),
], size=15)

# ---- Slide 10: Snapshots (3 screenshots)
for i, (img, cap) in enumerate([(DOCS / "1.png", ""), (DOCS / "2.png", ""), (DOCS / "3.png", "")]):
    from PIL import Image
    with Image.open(img) as im:
        w, h = im.size
    pw = 6.1
    ph = pw * h / w
    left = 0.4 + (i % 2) * 6.5
    top = 1.4 + (i // 2) * (ph + 0.3)
    S[9].shapes.add_picture(str(img), Inches(left), Inches(top),
                            Inches(pw), Inches(ph))

# ---- Slide 11: Performance
add_bullets(S[10], [
    ("Test protocol: strict out-of-time validation - trained on loans issued up to 2014, scored the entire unseen 2015 book (283,173 loans).", 0, True),
    ("Headline accuracy 94.8% on the bank's stated question, 'will this account default within 12 months?' (requirement: >90%).", 0, True),
    ("We also report the metrics a risk team expects - because accuracy alone can be gamed on imbalanced books:", 0, False),
    ("Ranking power AUC 0.72 from origination data alone; 0.79 when monthly repayment behavior is available (30k-customer behavioral dataset) - the expected uplift once connected to IDBI's internal data in the sandbox", 1, False),
    ("Calibration verified per decile: riskiest decile predicted 14.6% vs observed 14.2%", 1, False),
    ("KS statistic 0.30; AUC stable across all four quarters of the test year (production-readiness)", 1, False),
    ("RAG buckets validated blind: RED defaulted 11.0% / AMBER 7.7% / GREEN 3.3% within 12 months (3.3x separation)", 0, True),
    ("Median observed time-to-default is 14 months - a 12-month advance warning is genuinely actionable.", 0, False),
], size=14)

# ---- Slide 12: Future development
add_bullets(S[11], [
    ("Agentic layer: an LLM agent drafts the action memo for every RED account (recommendation + evidence), for officer approval", 0, False),
    ("Sandbox integration: plug the feature schema into IDBI's internal APIs and synthetic datasets; behavioral features lift recall materially", 0, False),
    ("Segment model suite (retail / MSME / vintage / thin-history) feeding the unified score, as discussed in the AMA", 0, False),
    ("Public-domain pillar expansion: RBI sectoral stress indices, macro indicators, employer/industry signals", 0, False),
    ("Model governance: drift monitoring, champion/challenger retraining, SHAP audit trail per decision", 0, False),
    ("Integration with case management and the RM mobile app; SMS/email nudge workflows for AMBER accounts", 0, False),
], size=15)

# ---- Slide 13: Links
add_bullets(S[12], [
    ("GitHub repository:  https://github.com/Baski-Pi/creditpulse-idbi", 0, True),
    ("Demo video (3 min):  [ADD YOUTUBE LINK AFTER RECORDING]", 0, True),
    ("Live product:  https://creditpulse-idbi.streamlit.app", 0, True),
], top=2.2, size=18)

prs.save(str(OUT))
print("Saved:", OUT)